#!/usr/bin/env python3
"""Build LinkedIn launch assets from Chat / Copilot / Voice / Integrations screenshots.

LinkedIn blurs PPTX uploads — it re-rasterizes slides at low resolution. Upload the
PDF or individual PNGs instead (1920×1080, 150 DPI).

Run capture first:
  python scripts/capture_mode_screenshots.py

Then:
  python scripts/build_linkedin_ppt.py

Outputs:
  exports/linkedin/Nexus_LinkedIn_Launch.pdf   ← upload this to LinkedIn
  exports/linkedin/slide-01-chat.png             ← or upload PNGs as a multi-image post
  exports/Nexus_LinkedIn_Launch.pptx             ← optional; avoid uploading to LinkedIn
"""

from __future__ import annotations

from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
SHOTS = ROOT / "exports" / "screenshots"
EXPORTS = ROOT / "exports"
LINKEDIN_DIR = EXPORTS / "linkedin"
OUT_PPTX = EXPORTS / "Nexus_LinkedIn_Launch.pptx"
OUT_PDF = LINKEDIN_DIR / "Nexus_LinkedIn_Launch.pdf"

# LinkedIn document/carousel — landscape fits UI screenshots; use portrait via --portrait
LINKEDIN_LANDSCAPE = (1920, 1080)
LINKEDIN_PORTRAIT = (1080, 1350)
LINKEDIN_DPI = 150
MIN_SOURCE_W = 1920
MIN_SOURCE_H = 1080
# Nexus dark canvas — matches app chrome so letterboxing is invisible
BG_RGB = (6, 8, 13)

SLIDES = [
    ("01-chat.png", "chat"),
    ("02-copilot.png", "copilot"),
    ("03-voice.png", "voice"),
    ("04-integrations.png", "integrations"),
]


def _sharpen(img):
    from PIL import ImageFilter

    return img.filter(ImageFilter.UnsharpMask(radius=1.2, percent=130, threshold=2))


def render_linkedin_slide(img_path: Path, size: tuple[int, int]):
    from PIL import Image

    canvas = Image.new("RGB", size, BG_RGB)
    with Image.open(img_path) as img:
        img = img.convert("RGB")
        if img.width < MIN_SOURCE_W or img.height < MIN_SOURCE_H:
            print(
                f"  WARN {img_path.name} is only {img.width}×{img.height} — "
                f"re-run capture_mode_screenshots.py for sharp output"
            )
        scale = min(size[0] / img.width, size[1] / img.height)
        # Never upscale — upscaling causes blur; only downscale from retina captures
        scale = min(scale, 1.0)
        new_w = max(1, int(img.width * scale))
        new_h = max(1, int(img.height * scale))
        if scale < 1.0:
            resized = _sharpen(img.resize((new_w, new_h), Image.Resampling.LANCZOS))
        else:
            resized = img
        x = (size[0] - new_w) // 2
        y = (size[1] - new_h) // 2
        canvas.paste(resized, (x, y))
    return canvas


def export_linkedin_images(size: tuple[int, int]) -> list[Path]:
    LINKEDIN_DIR.mkdir(parents=True, exist_ok=True)
    paths: list[Path] = []
    for idx, (filename, label) in enumerate(SLIDES, start=1):
        rendered = render_linkedin_slide(SHOTS / filename, size)
        out = LINKEDIN_DIR / f"slide-{idx:02d}-{label}.png"
        rendered.save(out, format="PNG", optimize=False, compress_level=1)
        paths.append(out)
        print(f"  PNG  {out} ({size[0]}×{size[1]})")
    return paths


def export_linkedin_pdf(png_paths: list[Path]) -> None:
    # img2pdf embeds PNGs losslessly — PIL PDF recompresses and softens UI text
    try:
        import img2pdf

        with open(OUT_PDF, "wb") as f:
            f.write(img2pdf.convert([str(p) for p in png_paths]))
        print(f"  PDF  {OUT_PDF} (lossless PNG embed — upload this to LinkedIn)")
    except ImportError:
        from PIL import Image

        images = [Image.open(p) for p in png_paths]
        try:
            images[0].save(
                OUT_PDF,
                "PDF",
                save_all=True,
                append_images=images[1:],
                resolution=LINKEDIN_DPI,
            )
            print(f"  PDF  {OUT_PDF} ({LINKEDIN_DPI} DPI — install img2pdf for sharper PDF)")
        finally:
            for img in images:
                img.close()


def build_pptx(size: tuple[int, int]) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    slide_w_in = size[0] / LINKEDIN_DPI
    slide_h_in = size[1] / LINKEDIN_DPI

    prs = Presentation()
    prs.slide_width = Inches(slide_w_in)
    prs.slide_height = Inches(slide_h_in)
    prs.core_properties.title = "Nexus — Chat · Copilot · Voice · Integrations"

    LINKEDIN_DIR.mkdir(parents=True, exist_ok=True)
    for idx, (filename, label) in enumerate(SLIDES, start=1):
        png = LINKEDIN_DIR / f"slide-{idx:02d}-{label}.png"
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(png), 0, 0, width=prs.slide_width, height=prs.slide_height)

    EXPORTS.mkdir(exist_ok=True)
    prs.save(str(OUT_PPTX))
    print(f"  PPTX {OUT_PPTX} (local use only — LinkedIn compresses PPTX uploads)")


def main() -> None:
    import argparse

    parser = argparse.ArgumentParser(description="Build LinkedIn launch PDF/PNGs from screenshots")
    parser.add_argument(
        "--portrait",
        action="store_true",
        help="Use 1080×1350 (4:5) slides — more mobile feed space on LinkedIn",
    )
    parser.add_argument("--no-pptx", action="store_true", help="Skip PPTX generation")
    args = parser.parse_args()

    missing = [name for name, _ in SLIDES if not (SHOTS / name).exists()]
    if missing:
        raise SystemExit(
            f"Missing screenshots: {missing}\n"
            "Run: python scripts/capture_mode_screenshots.py"
        )

    size = LINKEDIN_PORTRAIT if args.portrait else LINKEDIN_LANDSCAPE
    print(f"Rendering LinkedIn slides at {size[0]}×{size[1]} …")
    png_paths = export_linkedin_images(size)
    export_linkedin_pdf(png_paths)
    if not args.no_pptx:
        build_pptx(size)

    print(
        f"\nDone.\n"
        f"  Upload to LinkedIn: {OUT_PDF}\n"
        f"  Or multi-image post: {LINKEDIN_DIR}/slide-*.png\n"
        f"  Avoid uploading the .pptx — LinkedIn re-compresses it and text/UI will blur."
    )


if __name__ == "__main__":
    main()
