#!/usr/bin/env python3
"""
Build a production presentation deck (PPTX) for Nexus.

Outputs:
  - exports/Nexus_Production_Deep_Dive.pptx

This script is intentionally self-contained and safe to run repeatedly.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path


ROOT = Path(__file__).resolve().parent.parent
EXPORTS = ROOT / "exports"
SHOTS = EXPORTS / "screenshots"
EXPORTS.mkdir(exist_ok=True)


def _first_existing(paths: list[Path]) -> Path | None:
    for p in paths:
        if p.exists():
            return p
    return None


@dataclass(frozen=True)
class Img:
    path: Path
    caption: str


def main() -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    # Screenshots from exports/screenshots (run capture_mode_screenshots.py first).
    images: list[Img] = []

    login_img = _first_existing(
        [
            SHOTS / "00-login.png",
            Path.home()
            / ".cursor/projects/Users-disastershubz-voice-agents-voice-agents/assets/image-2b57a9b6-9666-4e93-9ab4-fc5db0f6141b.png",
        ]
    )
    if login_img:
        images.append(Img(login_img, "Auth gate: sign-in / create account"))

    ui_img = _first_existing(
        [
            SHOTS / "01-chat.png",
            Path.home()
            / ".cursor/projects/Users-disastershubz-voice-agents-voice-agents/assets/image-4baec4b8-06cf-4b6b-b561-bed0e6f26bf7.png",
            Path.home()
            / ".cursor/projects/Users-disastershubz-voice-agents-voice-agents/assets/image-a93bc782-0bd1-4172-8335-c6b6588e1ad2.png",
        ]
    )
    if ui_img:
        images.append(Img(ui_img, "Nexus console UI (chat / copilot / voice)"))

    integrations_img = _first_existing(
        [
            SHOTS / "04-integrations.png",
            Path.home()
            / ".cursor/projects/Users-disastershubz-voice-agents-voice-agents/assets/image-60843011-a065-479f-aa5f-46574e57ddc6.png",
        ]
    )
    if integrations_img:
        images.append(
            Img(
                integrations_img,
                "Integrations catalog — 62 native connectors (CRM, CCaaS, BI, HRIS, and more)",
            )
        )

    prs = Presentation()
    prs.core_properties.title = "Nexus — Production Deep Dive"
    prs.core_properties.subject = "Production-grade omnichannel AI agent platform"

    # ---- helpers ----
    def add_title_slide(title: str, subtitle: str) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

    def add_bullets(title: str, bullets: list[str], note: str | None = None) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.level = 0
        if note:
            slide.notes_slide.notes_text_frame.text = note

    def add_picture_slide(title: str, img: Img) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
        slide.shapes.title.text = title

        from PIL import Image as PILImage

        top = Inches(1.4)
        max_h = Inches(5.2)
        max_w = prs.slide_width - Inches(1.6)

        with PILImage.open(img.path) as pil_img:
            img_w, img_h = pil_img.size
        aspect = img_w / img_h

        pic_h = max_h
        pic_w = int(max_h * aspect)
        if pic_w > max_w:
            pic_w = max_w
            pic_h = int(max_w / aspect)

        left = int((prs.slide_width - pic_w) / 2)
        slide.shapes.add_picture(str(img.path), left, top, width=pic_w, height=pic_h)

        # Caption
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(6.7), Inches(12.0), Inches(0.5))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = img.caption
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(100, 116, 139)

        slide.notes_slide.notes_text_frame.text = (
            f"Screenshot: {img.path.name}\n\n"
            "Explain what the viewer is seeing, what problem this UI solves, and how it connects to the backend APIs."
        )

    # ---- deck content ----
    add_title_slide(
        "Nexus — Production Deep Dive",
        "Omnichannel AI agent platform (Chat · Copilot · Voice) with auth, observability, and operations.",
    )

    add_bullets(
        "Problem & Outcome",
        [
            "CX teams lose speed and accuracy when chat, voice, and internal tools are siloed.",
            "Nexus provides one orchestrator and one console for customer + agent workflows across channels.",
            "Outcomes: faster resolutions, consistent quality, and full operational visibility.",
            "Production baseline: HTTPS, JWT auth, rate limiting, audit logging, backups, systemd operation.",
        ],
        note="Frame the problem in contact-center language (AHT, CSAT, escalations, handoffs).",
    )

    add_bullets(
        "What We Do (Agentic CX)",
        [
            "AI Agents for Customers: resolve chat/voice inquiries end-to-end (intake → execute → follow-up)",
            "AI Agents for Frontline Teams: copilot assistance with context + next-best actions + drafts",
            "AI Agents for Operations: evaluation, insights, and continuous improvement loops",
            "Shared knowledge layer: RAG retrieval with persistence (Chroma), citations, and guardrails",
        ],
    )

    add_bullets(
        "Architecture (Conceptual)",
        [
            "Browser UI (static) → FastAPI `/api/v1/*`",
            "JWT auth (multi-tenant) gates all production endpoints",
            "Orchestrator routes requests to tools + KB retrieval",
            "Persistence: SQLite + Chroma on free tier (upgrade path: Postgres + Redis)",
        ],
    )

    add_bullets(
        "Production Hardening",
        [
            "APP_ENV=production (docs disabled, demo reset disabled)",
            "AUTH_REQUIRED=true (JWT required for APIs + WS)",
            "Registration closed after bootstrap (ALLOW_REGISTRATION=false)",
            "Admin-only integrations + audit logging",
            "Rate limiting middleware per endpoint group",
        ],
    )

    add_bullets(
        "Operations",
        [
            "systemd: `nexus.service` for reboot-safe uptime",
            "Daily backups: SQLite + Chroma + config (`nexus-backup.timer`)",
            "Edge hardening: Caddy TLS + security headers + CSP",
            "Monitoring: health endpoint + optional Sentry",
        ],
    )

    add_bullets(
        "Integrations (62 native connectors)",
        [
            "Public catalog at /integrations — searchable by category (CRM, ticketing, CCaaS, BI, HRIS, …)",
            "Each connector: encrypted vault credentials, status API, and proxy routes",
            "Examples: HubSpot, Salesforce, Zendesk, Freshdesk, PagerDuty, Snowflake, Epic, Workday",
            "iPaaS webhooks (n8n/Zapier) alongside native adapters for lifecycle events",
        ],
    )

    for i, img in enumerate(images):
        add_picture_slide(f"Product UI — Example {i+1}", img)

    add_bullets(
        "Why This Helps Industry",
        [
            "62 native integrations — CRM, ticketing, telephony, BI, and HRIS out of the box",
            "Composable platform: swap LLMs, connectors, and automation flows",
            "Clear security posture: auth-first, encrypted vault, audit log",
            "Operational maturity: backups, systemd, health checks, 215+ automated tests",
        ],
    )

    add_bullets(
        "Next Steps",
        [
            "Move to Oracle Always Free Ampere A1 (still $0) for Postgres + Redis",
            "Add admin UI for user provisioning (instead of API-only)",
            "Add WAF/bot protection (Cloudflare) and alerting dashboards",
        ],
    )

    out = EXPORTS / "Nexus_Production_Deep_Dive.pptx"
    prs.save(out)
    print(str(out))


if __name__ == "__main__":
    main()

