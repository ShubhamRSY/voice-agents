from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

EXPORTS = Path(__file__).resolve().parent.parent / "exports"
EXPORTS.mkdir(exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Professional palette
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFA, 0xFA, 0xFC)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
BLUE = RGBColor(0x2B, 0x6C, 0xE6)
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)
WHITE_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x10, 0xB9, 0x81)
TEAL = RGBColor(0x0E, 0xA5, 0x9E)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)

def set_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, t, w, h, color, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, left, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape

def add_text(slide, left, t, w, h, text, size=18, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb

def add_bullets(slide, left, t, w, h, items, size=16, color=DARK_TEXT, spacing=8, bold_first=False):
    tb = slide.shapes.add_textbox(left, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(spacing)
        if bold_first and i == 0:
            p.font.bold = True
    return tb

def slide_header(slide, number, subtitle, title, bg_color=NAVY):
    set_bg(slide, WHITE)
    # top bar
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.6), bg_color)
    add_text(slide, Inches(0.8), Inches(0.2), Inches(3), Inches(0.3),
             f"0{number}", size=11, color=WHITE_TEXT, bold=True)
    add_text(slide, Inches(0.8), Inches(0.45), Inches(8), Inches(0.3),
             subtitle.upper(), size=11, color=RGBColor(0x90, 0xAA, 0xCC))
    add_text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
             title, size=28, color=WHITE_TEXT, bold=True)
    # accent line under header
    add_rect(slide, Inches(0), Inches(1.6), Inches(13.333), Pt(3), BLUE)

TOTAL = 15

# ════════════════════ SLIDE 1: TITLE ════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, NAVY)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)

# accent bar
add_rect(s, Inches(0.8), Inches(3.0), Inches(1.5), Pt(5), BLUE)

add_text(s, Inches(0.8), Inches(1.8), Inches(11), Inches(1.2),
         "Nexus", size=64, color=WHITE_TEXT, bold=True)
add_text(s, Inches(0.8), Inches(3.4), Inches(11), Inches(0.6),
         "Purpose-built AI agents. One CX platform. (Open-source)", size=20, color=RGBColor(0x90, 0xAA, 0xCC))
add_text(s, Inches(0.8), Inches(4.3), Inches(11), Inches(0.4),
         "Chat  ·  Copilot  ·  Voice  —  One Orchestrator", size=16, color=RGBColor(0x70, 0x90, 0xB0))
add_text(s, Inches(0.8), Inches(5.5), Inches(11), Inches(0.4),
         "github.com/ShubhamRSY/voice-agents", size=12, color=RGBColor(0x50, 0x70, 0x90))

# ════════════════════ SLIDE 2: WHAT IS NEXUS ════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, 1, "Overview", "What Is Nexus?")

add_text(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.5),
         "Nexus is an open-source Agentic CX platform for customer experience teams. "
         "It solves the problem of fragmented support tools by providing a single AI-powered orchestrator "
         "that runs customer and agent workflows across chat, copilot, and voice.\n\n"
         "Purpose-built agents can resolve customer inquiries end-to-end, guide frontline teams in real time, "
         "and support operations with evaluation and insights — all grounded with a shared knowledge layer (RAG).\n\n"
         "Nexus ships with a built-in mock LLM so the platform runs without any API keys, making it fully functional "
         "for local development and testing out of the box.",
         size=15, color=DARK_TEXT)

# ════════════════════ SLIDE 3: CHANNELS ════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, 2, "Channels", "Three Channels, One Engine")

cards = [
    ("Chat", "Live AI conversation with SSE streaming,\nRAG citations, session history, and\nper-mode message filtering.", BLUE),
    ("Copilot", "Agent-assist mode — paste a customer\ntranscript and get an AI-generated\nreply in seconds.", TEAL),
    ("Voice", "PSTN telephony via Twilio, Amazon\nConnect, or any SIP trunk. Live STT →\nAI → TTS pipeline.", ORANGE),
]

for i, (title, desc, color) in enumerate(cards):
    x = Inches(0.8 + i * 4.1)
    y = Inches(2.2)
    card = add_rect(s, x, y, Inches(3.7), Inches(3.8), LIGHT_GRAY, radius=0.04)
    add_rect(s, x, y, Inches(3.7), Pt(4), color)
    add_text(s, x + Inches(0.3), y + Inches(0.3), Inches(3.2), Inches(0.4),
             title, size=20, color=color, bold=True)
    add_text(s, x + Inches(0.3), y + Inches(1.0), Inches(3.2), Inches(2.5),
             desc, size=14, color=GRAY)

# ════════════════════ SLIDE 4: ARCHITECTURE ════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, 3, "System Design", "Architecture")

flow = [
    "User / Console  →  REST API  →  Orchestrator",
    "",
    "  Session Manager — create, load, persist sessions",
    "  Prompt Builder — system prompt + chat history + RAG context",
    "  Agent Router — select model & agent by configuration",
    "",
    "  Flow:  Receive → Session → RAG → Prompt → LLM → Stream",
    "",
    "Services layer:  Feedback Engine · Vault (AES-256-GCM) · iPaaS Webhooks",
    "Infrastructure:  Twilio · Docker · PostgreSQL · Redis",
]
add_bullets(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(5),
            flow, size=15, color=DARK_TEXT, spacing=4)

# ════════════════════ SLIDE 5: MULTI-LLM ════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, 4, "LLM Layer", "Multi-LLM Support")

llms = [
    "OpenAI GPT-4o / GPT-4o-mini — default provider, production-ready",
    "Anthropic Claude 3.5 Sonnet & Haiku — alternative provider",
    "Google Gemini 2.0 Flash — via gemini_provider module",
    "Built-in Mock LLM — enables full development without API keys",
    "Per-agent model assignment via YAML configuration",
    "Pluggable provider interface — add any LLM in a single file",
]
add_bullets(s, Inches(0.8), Inches(2.2), Inches(11), Inches(4),
            llms, size=16, color=DARK_TEXT, spacing=10)

# ════════════════════ SLIDE 6: RAG ════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, 5, "Knowledge", "RAG Engine")

rag = [
    "Vector-based retrieval augmented generation",
    "Configurable chunk size & top-K retrieval parameters",
    "Source-grounded citations in every AI response",
    "Knowledge base → vector embeddings pipeline",
    "Per-agent knowledge base assignment in YAML",
    "Adjustable retrieval threshold per use case",
]
add_bullets(s, Inches(0.8), Inches(2.2), Inches(11), Inches(3.5),
            rag, size=16, color=DARK_TEXT, spacing=10)

add_rect(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5), LIGHT_BLUE, radius=0.04)
add_text(s, Inches(1.1), Inches(5.3), Inches(11), Inches(0.4),
         "RAG Flow", size=14, color=BLUE, bold=True)
add_text(s, Inches(1.1), Inches(5.7), Inches(11), Inches(0.7),
         "User Question  →  Embed  →  Vector Search  →  Retrieve Chunks  →  Augment Prompt  →  Generate Response with Citations",
         size=13, color=DARK_TEXT)

# ════════════════════ SLIDE 7: VOICE ════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, 6, "Telephony", "Voice & PSTN Integration")

voice = [
    "Twilio Programmable Voice — production-grade PSTN handling",
    "Amazon Connect — CCaaS integration for enterprise contact centers",
    "Generic SIP / CCaaS — compatible with any SIP trunk provider",
    "VAPI.ai & Retell AI — AI voice agent overlay support",
    "Built-in voice simulator — full testing without phone hardware",
    "Deepgram & AssemblyAI for speech-to-text",
    "ElevenLabs & PlayHT for text-to-speech",
    "Call status tracking and logs via REST API endpoints",
    "Twilio webhook handlers with status callbacks",
]
add_bullets(s, Inches(0.8), Inches(2.0), Inches(11), Inches(5),
            voice, size=14, color=DARK_TEXT, spacing=6)

# ════════════════════ SLIDE 8: FEEDBACK ════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, 7, "Optimization", "Feedback Engine")

fb = [
    "Users rate responses with thumbs up / down or CSAT score",
    "CSAT ratings dynamically tune agent temperature",
    "Adjusts RAG retrieval threshold based on feedback trends",
    "Sentiment tracking across sessions and channels",
    "All feedback logged for dashboard and analysis",
    "Enables continuous AI improvement without manual tuning",
]
add_bullets(s, Inches(0.8), Inches(2.2), Inches(11), Inches(3.5),
            fb, size=16, color=DARK_TEXT, spacing=10)

add_rect(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5), LIGHT_GRAY, radius=0.04)
add_text(s, Inches(1.1), Inches(5.3), Inches(11), Inches(0.4),
         "Feedback Loop", size=14, color=BLUE, bold=True)
add_text(s, Inches(1.1), Inches(5.7), Inches(11), Inches(0.7),
         "Rating Submitted  →  Temperature Adjusted  →  RAG Threshold Tuned  →  Next Response Improves",
         size=13, color=DARK_TEXT)

# ════════════════════ SLIDE 9: SECURITY ════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, 8, "Security", "Vault & Credential Management")

sec = [
    "AES-256-GCM encrypted vault — all credentials encrypted at rest",
    "REST API for credential CRUD operations",
    "JWT authentication with configurable expiry",
    "Pydantic input validation on every API endpoint",
    "Strict CORS origin whitelist",
    "Environment files never committed to version control",
    "No plaintext secrets in memory beyond request scope",
]
add_bullets(s, Inches(0.8), Inches(2.2), Inches(11), Inches(4),
            sec, size=16, color=DARK_TEXT, spacing=8)

# ════════════════════ SLIDE 10: CONSOLE UI ════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, 9, "Frontend", "Console UI")

ui = [
    "Dark mode interface with clean animations and micro-interactions",
    "Three-channel tabs: Chat, Copilot, and Voice",
    "Per-mode message filtering with sync toggle (combined or separated)",
    "Session sidebar — create, rename, switch, and clear sessions",
    "SSE streaming for word-by-word AI response display",
    "Typing indicator, streaming cursor, and entrance animations",
    "Welcome screen with per-mode placeholders",
    "Scrollable history with thin scrollbars and responsive layout",
]
add_bullets(s, Inches(0.8), Inches(2.2), Inches(11), Inches(4.5),
            ui, size=15, color=DARK_TEXT, spacing=7)

# ════════════════════ SLIDE 11: INTEGRATIONS ════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, 10, "Integration", "62 Native Connectors")

ip = [
    "62 native integrations — CRM, ticketing, CCaaS, telephony, BI, HRIS, knowledge, and project tools",
    "Public catalog at /integrations with search and category filters",
    "Encrypted vault (AES-256-GCM) + status API + 55+ provider proxy routes per connector",
    "Examples: HubSpot, Salesforce, Zendesk, Freshdesk, PagerDuty, Snowflake, Epic, Workday, Jira",
    "iPaaS webhooks for n8n/Zapier — lifecycle events alongside native adapters",
    "Nexus Cloud plans: Free / $29 Starter / $99 Growth with self-serve sign-up",
]
add_bullets(s, Inches(0.8), Inches(2.2), Inches(11), Inches(4.5),
            ip, size=15, color=DARK_TEXT, spacing=8)

# ════════════════════ SLIDE 12: AGENTS CONFIG ════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, 11, "Configuration", "Agent Definitions")

add_text(s, Inches(0.8), Inches(2.0), Inches(11), Inches(0.4),
         "Agents are defined in YAML files — no code changes required to add or modify agents:",
         size=14, color=GRAY)

# code block
add_rect(s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(3.5), LIGHT_GRAY, radius=0.04)
code = [
    "# config/agents/support.yaml",
    "name: acme_support",
    "display_name: ACME Support Agent",
    "model: gpt-4o-mini",
    "temperature: 0.7",
    "system_prompt: |",
    "  You are a helpful support agent for ACME Corp.",
    "  Answer questions based on the knowledge base.",
    "knowledge_base:",
    "  source: docs/knowledge_base/acme_support.md",
    "rag:",
    "  enabled: true",
    "  chunk_size: 512",
    "  top_k: 3",
]
add_bullets(s, Inches(1.0), Inches(2.7), Inches(11), Inches(3),
            code, size=11, color=GRAY, spacing=1)

# ════════════════════ SLIDE 13: DEPLOYMENT ════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, 12, "Operations", "Deployment Options")

deploy = [
    ("Docker", "docker compose -f deploy/docker/docker-compose.yml up"),
    ("Bare Metal", "uvicorn src.main:app --host 0.0.0.0 --port 8001"),
    ("CI/CD", "GitHub Actions pipeline — lint, test, e2e, build, deploy"),
]

for i, (method, cmd) in enumerate(deploy):
    y = Inches(2.3 + i * 1.3)
    add_rect(s, Inches(0.8), y, Inches(1.5), Inches(0.35), NAVY, radius=0.03)
    add_text(s, Inches(0.8), y, Inches(1.5), Inches(0.35),
             method, size=11, color=WHITE_TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.6), y + Pt(2), Inches(9), Inches(0.35),
             cmd, size=14, color=DARK_TEXT, font="Consolas")

# ════════════════════ SLIDE 14: TESTING ════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(s, 13, "Quality", "Testing & CI/CD")

test = [
    "215+ tests — unit, integration, E2E journeys, and live comprehensive E2E",
    "All 55 integration proxy routes verified (mock-safe, no 5xx without credentials)",
    "Mock LLM enables full test coverage without external API keys",
    "Ruff linter and mypy type checking enforced in CI",
    "Coverage reports output to reports/coverage/",
    "GitHub Actions workflow: lint → test → e2e → Docker build → deploy",
]
add_bullets(s, Inches(0.8), Inches(2.2), Inches(11), Inches(3.5),
            test, size=15, color=DARK_TEXT, spacing=8)

# badges row
for i, (name, color) in enumerate([("pytest", BLUE), ("ruff", TEAL), ("mypy", ORANGE)]):
    x = Inches(0.8 + i * 2.8)
    add_rect(s, x, Inches(5.5), Inches(2.2), Inches(0.4), color, radius=0.03)
    add_text(s, x, Inches(5.5), Inches(2.2), Inches(0.4),
             name, size=12, color=WHITE_TEXT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════ SLIDE 15: THANK YOU ════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, NAVY)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)
add_rect(s, Inches(5.5), Inches(3.5), Inches(2.5), Pt(5), BLUE)

add_text(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2),
         "Thank You", size=56, color=WHITE_TEXT, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(3.9), Inches(11.5), Inches(0.5),
         "Nexus — Open-Source Omnichannel AI Agent Platform", size=18, color=RGBColor(0x90, 0xAA, 0xCC), align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.4),
         "MIT License  ·  github.com/ShubhamRSY/voice-agents", size=13, color=RGBColor(0x60, 0x80, 0xA0), align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.4),
         "Chat  ·  Copilot  ·  Voice  —  One Orchestrator", size=11, color=RGBColor(0x40, 0x60, 0x80), align=PP_ALIGN.CENTER)

# ════════════════════ SAVE ════════════════════
path = EXPORTS / "Nexus_Overview.pptx"
prs.save(path)
print(f"Saved {path} — 15 slides")
