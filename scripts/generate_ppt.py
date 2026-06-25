from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import math

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Modern Palette ──
BG = RGBColor(0x0B, 0x0B, 0x1A)        # deep navy-black
BG2 = RGBColor(0x11, 0x11, 0x28)        # slightly lighter
ACCENT = RGBColor(0x6C, 0x63, 0xFF)     # purple-primary
ACCENT2 = RGBColor(0x00, 0xD4, 0xAA)    # teal-secondary
ACCENT3 = RGBColor(0xFF, 0x6B, 0x6B)    # coral
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xB0, 0xB0, 0xC8)
MUTED = RGBColor(0x6B, 0x6B, 0x8A)

def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape

def add_circle(slide, l, t, size, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, l, t, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb

def add_bullets(slide, l, t, w, h, items, size=16, color=LIGHT, spacing=12):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(spacing)
    return tb

def accent_bar(slide, y):
    add_rect(slide, Inches(0.8), y, Inches(1.2), Pt(4), ACCENT2)

def section_label(slide, text, y=Inches(0.7)):
    add_text(slide, Inches(0.8), Inches(0.3), Inches(6), Inches(0.4),
             text.upper(), size=12, color=ACCENT2, bold=True)
    accent_bar(slide, Inches(0.7))

def section_title(slide, text, y=Inches(1.1)):
    add_text(slide, Inches(0.8), y, Inches(11), Inches(1),
             text, size=40, color=WHITE, bold=True)

def page_number(slide, num, total):
    add_text(slide, Inches(12), Inches(7.0), Inches(1), Inches(0.3),
             f"{num:02d} / {total:02d}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)

def add_badge(slide, text, l, t, color=ACCENT):
    s = add_rect(slide, l, t, Inches(2.4), Inches(0.35), color, radius=0.05)
    tf = s.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    return s

TOTAL = 16

# ════════════════════════════ SLIDE 1: TITLE ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

# Decorative circles
add_circle(s, Inches(9.5), Inches(-1.5), Inches(5), ACCENT)
add_circle(s, Inches(10.5), Inches(4), Inches(2.5), ACCENT2)
add_circle(s, Inches(-1), Inches(5), Inches(3), ACCENT)

# Central content
add_text(s, Inches(1), Inches(1.5), Inches(8), Inches(1.2),
         "NEXUS", size=80, color=WHITE, bold=True)
add_rect(s, Inches(1), Inches(2.8), Inches(2), Pt(5), ACCENT2)
add_text(s, Inches(1), Inches(3.2), Inches(8), Inches(0.6),
         "Open-Source Omnichannel AI Agent Platform", size=24, color=ACCENT)
add_text(s, Inches(1), Inches(4.2), Inches(8), Inches(0.5),
         "Chat  ·  Copilot  ·  Voice  —  One Orchestrator", size=16, color=LIGHT)
add_text(s, Inches(1), Inches(5.5), Inches(8), Inches(0.4),
         "github.com/ShubhamRSY/voice-agents", size=12, color=MUTED)
page_number(s, 1, TOTAL)

# ════════════════════════════ SLIDE 2: WHAT IS NEXUS ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_circle(s, Inches(-0.8), Inches(5.5), Inches(2.5), ACCENT)

section_label(s, "Overview")
section_title(s, "What Is Nexus?")

add_text(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(3.5),
         "Customer support teams lose context switching between channels — chat, phone, email, and "
         "internal tools each live in separate silos. Nexus solves this with a single AI-powered "
         "orchestrator that handles every conversation from one runtime, with one knowledge base, "
         "and one feedback loop that improves responses over time.\n\n"
         "No API key is required to start — the platform ships with a built-in mock LLM so the "
         "entire console, voice simulator, and all 158+ unit tests work immediately out of the box.",
         size=16, color=LIGHT)
page_number(s, 2, TOTAL)

# ════════════════════════════ SLIDE 3: THREE CHANNELS ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

section_label(s, "Channels")
section_title(s, "Three Channels, One Engine")

cards = [
    ("💬  Chat", "Live AI conversation with SSE streaming,\nRAG citations, session history, and\nper-mode message filtering.", ACCENT),
    ("🤝  Copilot", "Agent-assist mode — paste a customer\ntranscript, get an AI-generated reply\ninstantly.", ACCENT2),
    ("📞  Voice", "PSTN telephony via Twilio, Amazon\nConnect, or any SIP trunk. Live\nSTT → AI → TTS pipeline.", ACCENT3),
]

for i, (title, desc, color) in enumerate(cards):
    x = Inches(0.8 + i * 4.1)
    y = Inches(2.2)
    card = add_rect(s, x, y, Inches(3.7), Inches(4.2), BG2, radius=0.05)
    # top accent
    add_rect(s, x, y, Inches(3.7), Pt(5), color)
    add_text(s, x + Inches(0.3), y + Inches(0.4), Inches(3.2), Inches(0.5),
             title, size=22, color=color, bold=True)
    add_text(s, x + Inches(0.3), y + Inches(1.2), Inches(3.2), Inches(2.5),
             desc, size=14, color=LIGHT)

page_number(s, 3, TOTAL)

# ════════════════════════════ SLIDE 4: ARCHITECTURE ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

section_label(s, "System Design")
section_title(s, "Architecture")

add_circle(s, Inches(11), Inches(5), Inches(2), ACCENT2)
add_circle(s, Inches(-0.5), Inches(-0.5), Inches(1.5), ACCENT)

flow = [
    "User / Console  ──→  REST API  ──→  Orchestrator",
    "",
    "  │  Session Manager — create, load, persist",
    "  │  Prompt Builder — system prompt + history + context",
    "  │  Agent Router — model & agent selection",
    "  │",
    "  ├──→  RAG Engine  —  vector retrieval + citations",
    "  ├──→  LLM Layer  —  OpenAI · Anthropic · Gemini · Mock",
    "  └──→  SSE Stream  ──→  Console UI",
    "",
    "Services:  Feedback Engine  ·  Vault (AES-256-GCM)  ·  iPaaS Webhooks",
    "Infra:    Twilio  ·  Redis  ·  PostgreSQL  ·  Docker",
]
add_bullets(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.5),
            flow, size=15, color=LIGHT, spacing=4)
page_number(s, 4, TOTAL)

# ════════════════════════════ SLIDE 5: MULTI-LLM ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

section_label(s, "LLM Layer")
section_title(s, "Multi-LLM Support")

add_circle(s, Inches(10.5), Inches(1.5), Inches(1.8), ACCENT3)

llms = [
    "OpenAI GPT-4o / GPT-4o-mini — default provider",
    "Anthropic Claude 3.5 — Sonnet & Haiku",
    "Google Gemini 2.0 Flash — via gemini_provider",
    "Built-in Mock LLM — full dev parity, no API key needed",
    "Per-agent model config in YAML — switch at deploy time",
    "Pluggable provider interface — add any LLM in one file",
]
add_bullets(s, Inches(0.8), Inches(2.2), Inches(11), Inches(4),
            llms, size=17, color=LIGHT, spacing=10)
page_number(s, 5, TOTAL)

# ════════════════════════════ SLIDE 6: RAG ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

section_label(s, "Knowledge")
section_title(s, "RAG Engine")

rag = [
    "Vector-based retrieval augmented generation",
    "Configurable chunk size & top-K retrieval",
    "Source-grounded citations in every response",
    "Knowledge base → vector embeddings pipeline",
    "Per-agent KB config in YAML agent definitions",
    "Adjustable retrieval threshold per use case",
]
add_bullets(s, Inches(0.8), Inches(2.2), Inches(11), Inches(4),
            rag, size=17, color=LIGHT, spacing=10)

add_rect(s, Inches(9), Inches(3.5), Inches(3.5), Inches(2.5), BG2, radius=0.05)
add_text(s, Inches(9.3), Inches(3.8), Inches(3), Inches(0.4),
         "RAG Flow", size=14, color=ACCENT2, bold=True)
rag_flow = [
    "Query → Embed → Retrieve → Augment → Generate",
    "",
    "User question chunked,",
    "matched against KB,",
    "injected into prompt",
    "with source citations",
]
add_bullets(s, Inches(9.3), Inches(4.3), Inches(3), Inches(2),
            rag_flow, size=11, color=LIGHT, spacing=4)

page_number(s, 6, TOTAL)

# ════════════════════════════ SLIDE 7: VOICE ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

section_label(s, "Telephony")
section_title(s, "Voice & PSTN")

add_circle(s, Inches(-0.5), Inches(4), Inches(2), ACCENT3)

voice = [
    "Twilio Programmable Voice — production PSTN",
    "Amazon Connect — CCaaS integration",
    "Generic SIP / CCaaS — any SIP trunk supported",
    "VAPI.ai & Retell AI — AI voice agent overlays",
    "Built-in simulator — full voice testing without hardware",
    "Deepgram & AssemblyAI for STT",
    "ElevenLabs & PlayHT for TTS",
    "Call status tracking & logs via REST API",
    "Twilio webhooks + status callbacks",
]
add_bullets(s, Inches(0.8), Inches(2.2), Inches(11), Inches(4.5),
            voice, size=16, color=LIGHT, spacing=6)
page_number(s, 7, TOTAL)

# ════════════════════════════ SLIDE 8: FEEDBACK ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

section_label(s, "Optimization")
section_title(s, "Feedback Engine")

add_circle(s, Inches(11), Inches(0.5), Inches(1.5), ACCENT2)

fb = [
    "User rates responses — 👍 / 👎 / CSAT score",
    "CSAT dynamically tunes agent temperature",
    "Adjusts RAG retrieval threshold based on feedback data",
    "Sentiment tracking across sessions & channels",
    "All feedback logged for dashboard & analysis",
    "Enables continuous improvement without manual tuning",
]
add_bullets(s, Inches(0.8), Inches(2.2), Inches(11), Inches(3.5),
            fb, size=17, color=LIGHT, spacing=8)

# feedback loop box
add_rect(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5), BG2, radius=0.05)
add_text(s, Inches(1.1), Inches(5.4), Inches(11), Inches(0.4),
         "Feedback Loop", size=14, color=ACCENT2, bold=True)
add_text(s, Inches(1.1), Inches(5.8), Inches(11), Inches(0.7),
         "Rating Submitted  →  Temperature Adjusted  →  RAG Threshold Tuned  →  Next Response Improves",
         size=14, color=LIGHT)
page_number(s, 8, TOTAL)

# ════════════════════════════ SLIDE 9: SECURITY ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

section_label(s, "Security")
section_title(s, "Vault & Credential Management")

sec = [
    "AES-256-GCM encrypted vault — all credentials at rest",
    "REST API for CRUD on integration credentials",
    "JWT authentication with configurable expiry",
    "Pydantic input validation on every endpoint",
    "Strict CORS origin whitelist",
    ".env never committed — all secrets externalized",
    "No plaintext keys stored in memory beyond request scope",
]
add_bullets(s, Inches(0.8), Inches(2.2), Inches(11), Inches(4),
            sec, size=16, color=LIGHT, spacing=8)

add_rect(s, Inches(9.5), Inches(3.5), Inches(3), Inches(2.8), BG2, radius=0.05)
add_text(s, Inches(9.8), Inches(3.7), Inches(2.5), Inches(0.4),
         "Encryption Flow", size=13, color=ACCENT2, bold=True)
enc = [
    "Key in .env",
    "  ↓",
    "AES-256-GCM",
    "  ↓",
    "Ciphertext → DB",
    "  ↓",
    "Decrypt on read",
]
add_bullets(s, Inches(9.8), Inches(4.2), Inches(2.5), Inches(2),
            enc, size=11, color=LIGHT, spacing=2)
page_number(s, 9, TOTAL)

# ════════════════════════════ SLIDE 10: CONSOLE UI ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

section_label(s, "Frontend")
section_title(s, "Console UI")

add_circle(s, Inches(11.5), Inches(4.5), Inches(2), ACCENT2)

ui = [
    "Dark mode with premium animations & micro-interactions",
    "Three-channel tabs: Chat / Copilot / Voice",
    "Per-mode message filtering + sync toggle (🔗 / ⊘)",
    "Session sidebar — create, rename, switch, clear",
    "SSE streaming — word-by-word token display",
    "Typing indicator, streaming cursor, entrance animations",
    "Welcome screen with per-mode placeholders",
    "Scrollable history, button glow effects, thin scrollbars",
]
add_bullets(s, Inches(0.8), Inches(2.2), Inches(11), Inches(4.5),
            ui, size=16, color=LIGHT, spacing=6)
page_number(s, 10, TOTAL)

# ════════════════════════════ SLIDE 11: IPAAS ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

section_label(s, "Integration")
section_title(s, "iPaaS Webhooks")

add_circle(s, Inches(-0.3), Inches(0.5), Inches(1.8), ACCENT3)

ip = [
    "Lifecycle event webhooks for n8n, Zapier, and custom flows",
    "Events: session.created · message.completed · feedback.submitted",
    "Async HTTP dispatch via httpx — non-blocking",
    "Event history endpoint for debugging & replay",
    "Connect to CRMs: Salesforce, Zendesk, ServiceNow",
    "Notifications: Slack, email, and more",
    "Template workflows included for n8n & Zapier",
]
add_bullets(s, Inches(0.8), Inches(2.2), Inches(11), Inches(4),
            ip, size=16, color=LIGHT, spacing=8)
page_number(s, 11, TOTAL)

# ════════════════════════════ SLIDE 12: AGENTS CONFIG ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

section_label(s, "Configuration")
section_title(s, "Agent Definitions")

add_text(s, Inches(0.8), Inches(2.0), Inches(11), Inches(0.5),
         "Agents are defined in YAML — no code changes needed:", size=16, color=LIGHT)

# code block visual
code_bg = add_rect(s, Inches(0.8), Inches(2.6), Inches(11), Inches(3.2), RGBColor(0x15, 0x15, 0x25), radius=0.05)
code = [
    "# config/agents/support.yaml",
    "name: acme_support",
    "display_name: ACME Support Agent",
    "model: gpt-4o-mini",
    "temperature: 0.7",
    "system_prompt: |",
    "  You are a helpful support agent for ACME Corp.",
    "  Answer questions based on the knowledge base.",
    "knowledge_base:",
    "  source: docs/knowledge_base/acme_support.md",
    "rag:",
    "  enabled: true",
    "  chunk_size: 512",
    "  top_k: 3",
]
add_bullets(s, Inches(1.1), Inches(2.9), Inches(10.5), Inches(2.8),
            code, size=13, color=RGBColor(0xA0, 0xCC, 0x88), spacing=1)
page_number(s, 12, TOTAL)

# ════════════════════════════ SLIDE 13: DEPLOYMENT ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

section_label(s, "Operations")
section_title(s, "Deployment Options")

add_circle(s, Inches(10.5), Inches(5), Inches(2), ACCENT)

deploy_items = [
    ("Docker", "docker compose -f deploy/docker/docker-compose.yml up"),
    ("Bare Metal", "uvicorn src.main:app --host 0.0.0.0 --port 8001"),
    ("CI/CD", "GitHub Actions — lint, test, e2e, build, deploy"),
]

for i, (method, cmd) in enumerate(deploy_items):
    y = Inches(2.3 + i * 1.3)
    badge = add_rect(s, Inches(0.8), y, Inches(1.6), Inches(0.35), ACCENT, radius=0.04)
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = method
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    add_text(s, Inches(2.7), y, Inches(9), Inches(0.35),
             cmd, size=14, color=ACCENT2, font="Consolas")

page_number(s, 13, TOTAL)

# ════════════════════════════ SLIDE 14: TESTING ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

section_label(s, "Quality")
section_title(s, "Testing & CI/CD")

add_circle(s, Inches(-0.3), Inches(5.5), Inches(2), ACCENT2)

test = [
    "158+ unit tests — routers, services, providers, vault",
    "33 comprehensive E2E tests — real user flow simulation",
    "Mock LLM provider — full coverage without API keys",
    "Ruff linter + mypy type checking in CI pipeline",
    "Coverage reports → reports/coverage/",
    "GitHub Actions: lint → test → e2e → docker build → deploy",
]
add_bullets(s, Inches(0.8), Inches(2.2), Inches(11), Inches(4),
            test, size=16, color=LIGHT, spacing=8)

# badge row
add_badge(s, "pytest", Inches(0.8), Inches(5.5), ACCENT)
add_badge(s, "ruff", Inches(3.5), Inches(5.5), ACCENT2)
add_badge(s, "mypy", Inches(6.2), Inches(5.5), ACCENT3)
page_number(s, 14, TOTAL)

# ════════════════════════════ SLIDE 15: STACK ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

section_label(s, "Technology")
section_title(s, "Tech Stack")

stack = [
    ("Runtime", "Python 3.11+ · FastAPI · Uvicorn · SSE-Starlette"),
    ("LLM", "OpenAI · Anthropic · Google Gemini · Mock Provider"),
    ("Voice", "Twilio · Amazon Connect · Deepgram · ElevenLabs"),
    ("Data", "RAG (Vector) · PostgreSQL · Redis · Pydantic"),
    ("Frontend", "Vanilla JS · CSS · HTML · SSE Streaming"),
    ("Infra", "Docker · GitHub Actions · n8n · Zapier"),
    ("Security", "AES-256-GCM · JWT · CORS · Pydantic Validation"),
]

for i, (layer, desc) in enumerate(stack):
    y = Inches(2.2 + i * 0.7)
    add_rect(s, Inches(0.8), y, Inches(1.8), Inches(0.45), ACCENT, radius=0.03)
    add_text(s, Inches(0.8), y, Inches(1.8), Inches(0.45),
             layer, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2.9), y + Pt(3), Inches(9), Inches(0.4),
             desc, size=14, color=LIGHT)

page_number(s, 15, TOTAL)

# ════════════════════════════ SLIDE 16: THANK YOU ════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

# decorative
add_circle(s, Inches(-1), Inches(-1), Inches(3), ACCENT)
add_circle(s, Inches(11), Inches(5.5), Inches(2.5), ACCENT2)
add_circle(s, Inches(8), Inches(-0.5), Inches(1.5), ACCENT3)

add_text(s, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
         "Thank You", size=72, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_rect(s, Inches(5.5), Inches(3.3), Inches(2.5), Pt(5), ACCENT2)
add_text(s, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
         "Nexus — Open-Source Omnichannel AI Agent Platform", size=20, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.8), Inches(11), Inches(0.5),
         "MIT License  ·  github.com/ShubhamRSY/voice-agents", size=14, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
         "Chat  ·  Copilot  ·  Voice  —  One Orchestrator", size=12, color=MUTED, align=PP_ALIGN.CENTER)

page_number(s, 16, TOTAL)

# ════════════════════════════ SAVE ════════════════════════════
path = "Nexus_Overview.pptx"
prs.save(path)
print(f"Saved {path} — 16 slides")
